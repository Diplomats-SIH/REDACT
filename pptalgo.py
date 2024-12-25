from PyQt5.QtWidgets import (
    QApplication, QMainWindow, QVBoxLayout, QWidget, QLabel, QPushButton, QFileDialog, QSlider, QRadioButton, 
    QButtonGroup, QMessageBox, QHBoxLayout, QFrame, QProgressBar
)
from PyQt5.QtCore import Qt
from PyQt5.QtGui import QIcon
import sys
import os
from pptx import Presentation
import spacy
import re
from faker import Faker

# Initialize SpaCy and Faker
nlp = spacy.load("en_core_web_sm")
fake = Faker()

# Define redaction levels
REDACTION_LEVELS = {
    0: [],
    25: ["ORG", "EMAIL", "PHONE"],
    50: ["ORG", "EMAIL", "PHONE", "MONEY", "IP"],
    75: ["ORG", "EMAIL", "PHONE", "MONEY", "IP", "DATE", "TIME"],
    100: ["ORG", "EMAIL", "PHONE", "MONEY", "IP", "DATE", "TIME", "ADDRESS", "PERSON"]
}

# Function to extract text from PowerPoint
def extract_text_from_ppt(ppt_path):
    prs = Presentation(ppt_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text.append(shape.text)
    return "\n".join(text)

# Function to identify sensitive entities
def detect_entities(text, redaction_labels):
    doc = nlp(text)
    entities = []
    for ent in doc.ents:
        if ent.label_ in redaction_labels:
            entities.append((ent.text, ent.label_))
    # Additional regex-based detection for MONEY, TIME, and sensitive data
    if "MONEY" in redaction_labels:
        money_matches = re.findall(r"\$\d+(\.\d{2})?", text)
        entities.extend([(match, "MONEY") for match in money_matches])
    if "TIME" in redaction_labels:
        time_matches = re.findall(r"\b\d{1,2}:\d{2}(?:\s?[APap][Mm])?\b", text)
        entities.extend([(match, "TIME") for match in time_matches])
    if "IP" in redaction_labels:
        ip_matches = re.findall(r"\b(?:\d{1,3}\.){3}\d{1,3}\b", text)
        entities.extend([(match, "IP") for match in ip_matches])
    return entities

# Function to generate synthetic replacements for entities
def generate_synthetic_data(label):
    if label == "PERSON":
        return fake.name()
    elif label == "ORG":
        return fake.company()
    elif label == "EMAIL":
        return fake.email()
    elif label == "PHONE":
        return fake.phone_number()
    elif label == "MONEY":
        return fake.pricetag()
    elif label == "IP":
        return fake.ipv4()
    elif label == "DATE":
        return fake.date()
    elif label == "TIME":
        return fake.time()
    elif label == "ADDRESS":
        return fake.address()
    else:
        return "[SYNTHETIC DATA]"

# Function to apply redaction or synthetic data
def apply_redaction(text, entities, style="blackout"):
    for entity, label in entities:
        if style == "blackout":
            redacted = "█" * len(entity)
        elif style == "blur":
            redacted = "*" * len(entity)
        elif style == "synthetic":
            redacted = generate_synthetic_data(label)
        else:
            redacted = "[REDACTED]"
        text = re.sub(re.escape(entity), redacted, text)
    return text

# Function to redact text in PowerPoint
def redact_ppt(ppt_path, output_path, level, style="blackout"):
    prs = Presentation(ppt_path)
    redaction_labels = REDACTION_LEVELS.get(level, [])
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                original_text = shape.text
                entities = detect_entities(original_text, redaction_labels)
                redacted_text = apply_redaction(original_text, entities, style)
                shape.text = redacted_text
    
    prs.save(output_path)

class RedactApp(QMainWindow):
    def __init__(self):
        super().__init__()

        self.setWindowTitle("RE-DACT - PowerPoint Redactor")
        self.setGeometry(100, 100, 800, 600)
        self.ppt_path = None
        self.redacted_ppt_path = None

        self.setStyleSheet("""
            QMainWindow { background-color: #2C2F33; }
            QLabel { color: white; font-size: 16px; }
            QPushButton {
                font-size: 16px;
                padding: 10px;
                border-radius: 8px;
                background-color: #7289DA;
                color: white;
                border: none;
            }
            QPushButton:hover { background-color: #5A73BF; }
            QSlider::groove:horizontal {
                height: 8px; background: #444; border-radius: 4px;
            }
            QSlider::handle:horizontal {
                background: #7289DA; width: 18px; margin: -4px 0; border-radius: 9px;
            }
            QRadioButton { color: white; font-size: 16px; }
            QFrame {
                border: 1px solid #444;
                border-radius: 10px;
                background-color: #23272A;
                padding: 10px;
            }
        """)

        main_widget = QWidget()
        layout = QVBoxLayout()

        # File Selection Frame
        file_frame = QFrame()
        file_layout = QVBoxLayout()
        self.file_button = QPushButton("Choose PowerPoint File")
        self.file_button.setIcon(QIcon("file-icon.png"))  # Replace with the path to an icon file
        self.file_button.clicked.connect(self.choose_file)
        file_layout.addWidget(self.file_button)

        self.file_label = QLabel("No file selected")
        file_layout.addWidget(self.file_label)
        file_frame.setLayout(file_layout)
        layout.addWidget(file_frame)

        # Redaction Level Frame
        level_frame = QFrame()
        level_layout = QVBoxLayout()

        self.level_label = QLabel("Redaction Level")
        level_layout.addWidget(self.level_label)

        self.slider = QSlider(Qt.Horizontal)
        self.slider.setMinimum(0)
        self.slider.setMaximum(100)
        self.slider.setValue(50)
        level_layout.addWidget(self.slider)

        level_frame.setLayout(level_layout)
        layout.addWidget(level_frame)

        # Redaction Style Frame
        style_frame = QFrame()
        style_layout = QVBoxLayout()

        self.style_label = QLabel("Redaction Style")
        style_layout.addWidget(self.style_label)

        self.style_group = QButtonGroup(self)
        self.blur_button = QRadioButton("Blur Out")
        self.blackout_button = QRadioButton("Black Out")
        self.synthetic_button = QRadioButton("Synthetic Data Replacement")

        self.style_group.addButton(self.blur_button)
        self.style_group.addButton(self.blackout_button)
        self.style_group.addButton(self.synthetic_button)
        style_layout.addWidget(self.blur_button)
        style_layout.addWidget(self.blackout_button)
        style_layout.addWidget(self.synthetic_button)

        style_frame.setLayout(style_layout)
        layout.addWidget(style_frame)

        # Redact Button
        self.redact_button = QPushButton("Redact")
        self.redact_button.clicked.connect(self.redact_ppt)
        layout.addWidget(self.redact_button)

        # Progress Bar
        self.progress_bar = QProgressBar()
        self.progress_bar.setValue(0)
        self.progress_bar.setVisible(False)
        layout.addWidget(self.progress_bar)

        # Download Button
        self.download_button = QPushButton("Download Redacted File")
        self.download_button.setIcon(QIcon("download-icon.png"))  # Replace with the path to an icon file
        self.download_button.setEnabled(False)
        self.download_button.clicked.connect(self.download_redacted_file)
        layout.addWidget(self.download_button)

        main_widget.setLayout(layout)
        self.setCentralWidget(main_widget)

    def choose_file(self):
        options = QFileDialog.Options()
        file_name, _ = QFileDialog.getOpenFileName(self, "Choose PowerPoint File", "", "PowerPoint Files (*.pptx)", options=options)
        if file_name:
            self.ppt_path = file_name
            self.file_label.setText(file_name)

    def redact_ppt(self):
        if not self.ppt_path:
            QMessageBox.warning(self, "Error", "Please select a PowerPoint file first.")
            return

        redact_level = self.slider.value()
        if self.blur_button.isChecked():
            action = 'blur'
        elif self.blackout_button.isChecked():
            action = 'blackout'
        elif self.synthetic_button.isChecked():
            action = 'synthetic'
        else:
            QMessageBox.warning(self, "Error", "Please select a redaction style.")
            return

        self.progress_bar.setVisible(True)
        self.progress_bar.setValue(0)

        try:
            output_path = os.path.splitext(self.ppt_path)[0] + "_redacted.pptx"
            redact_ppt(self.ppt_path, output_path, redact_level, style=action)
            self.redacted_ppt_path = output_path
            self.progress_bar.setValue(100)
            self.download_button.setEnabled(True)
            QMessageBox.information(self, "Success", f"Redaction completed! File saved at: {output_path}")
        except Exception as e:
            QMessageBox.warning(self, "Error", f"An error occurred: {str(e)}")
        finally:
            self.progress_bar.setVisible(False)

    def download_redacted_file(self):
        if self.redacted_ppt_path:
            options = QFileDialog.Options()
            save_path, _ = QFileDialog.getSaveFileName(self, "Save Redacted File", self.redacted_ppt_path, "PowerPoint Files (*.pptx)", options=options)
            if save_path:
                os.rename(self.redacted_ppt_path, save_path)
                self.redacted_ppt_path = save_path
                QMessageBox.information(self, "Success", f"File downloaded successfully to {save_path}")

if __name__ == "__main__":
    app = QApplication(sys.argv)
    window = RedactApp()
    window.show()
    sys.exit(app.exec_())
